import asyncio
import base64
import fitz
import hashlib
import os
from typing import Optional
from urllib.parse import urlparse
from fastapi import HTTPException
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import WebBaseLoader
from langchain_openai import OpenAIEmbeddings
from openai import AuthenticationError, OpenAI, OpenAIError, RateLimitError
from pinecone.exceptions import PineconeApiException
from app.config import (
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    EMBEDDING_DIMENSIONS,
    EMBEDDING_MODEL,
    MIN_IMAGE_BYTES,
    TOP_K,
    VISION_MODEL,
)
from app.db.vector_store import get_index

MIN_TEXT_CHARS = 100
# Chars-per-page threshold below which a PDF is treated as scanned
_SCANNED_CHARS_PER_PAGE = 50


def get_openai_client():
    return OpenAI()


def chunk_vector_id(collection_id: str, source: str, chunk_index: int) -> str:
    """Deterministic Pinecone vector ID derived from content coordinates.
    Ensures re-running a job never creates duplicate vectors."""
    key = f"{collection_id}|{source}|{chunk_index}"
    return hashlib.sha256(key.encode()).hexdigest()

async def ingest_pdf(file, namespace: str = ""):
    pdf_bytes = await file.read()
    file_name = os.path.basename(file.filename)
    return await ingest_pdf_from_bytes(pdf_bytes, file_name, namespace=namespace)


async def ingest_pdf_from_bytes(pdf_bytes: bytes, file_name: str, namespace: str = ""):
    text = extract_pdf_text(pdf_bytes)
    chunks = chunk_text(text)
    store_chunks(chunks, file_name, namespace=namespace)
    return {
        "message": "Document ingested successfully",
        "chunks": len(chunks)
    }


async def ingest_url(url: str, namespace: str = ""):
    loader = WebBaseLoader(url)
    docs = await asyncio.get_event_loop().run_in_executor(None, loader.load)
    text = "\n\n".join(doc.page_content for doc in docs)

    parsed = urlparse(url)
    source_name = parsed.netloc + parsed.path

    chunks = chunk_text(text)
    store_chunks(chunks, source_name, namespace=namespace)

    return {
        "message": "URL ingested successfully",
        "chunks": len(chunks)
    }


def extract_pdf_text(pdf_bytes):
    doc = fitz.open(
        stream=pdf_bytes,
        filetype="pdf"
    )

    text = ""

    for page in doc:
        text += page.get_text()

    return text

def _is_scanned_pdf(doc: fitz.Document) -> bool:
    """Returns True when avg text per page is below the scanned-PDF threshold."""
    if len(doc) == 0:
        return False
    total = sum(len(page.get_text().strip()) for page in doc)
    return (total / len(doc)) < _SCANNED_CHARS_PER_PAGE


def _table_to_markdown(rows: list[list]) -> str:
    """Serialize a PyMuPDF table (list of rows) to a GFM markdown table string."""
    if not rows:
        return ""
    clean = [[str(c).strip() if c is not None else "" for c in row] for row in rows]
    clean = [r for r in clean if any(c for c in r)]
    if not clean:
        return ""
    col_count = max(len(r) for r in clean)
    clean = [r + [""] * (col_count - len(r)) for r in clean]
    lines = [
        "| " + " | ".join(clean[0]) + " |",
        "| " + " | ".join(["---"] * col_count) + " |",
    ]
    for row in clean[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _page_text_excluding_tables(page: fitz.Page, table_rects: list) -> str:
    """Return page text, skipping blocks that overlap >50% with a detected table region."""
    if not table_rects:
        return page.get_text()
    parts = []
    for block in page.get_text("blocks"):
        if block[6] != 0:
            continue
        brect = fitz.Rect(block[:4])
        block_area = brect.get_area()
        if block_area == 0:
            continue
        if any((brect & tr).get_area() / block_area > 0.5 for tr in table_rects):
            continue
        parts.append(block[4])
    return "\n".join(parts)


def _describe_image_vision_sync(image_bytes: bytes, ext: str, client: OpenAI) -> str:
    """Call the vision model synchronously and return a description string."""
    mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
    b64 = base64.b64encode(image_bytes).decode()
    resp = client.chat.completions.create(
        model=VISION_MODEL,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                {"type": "text", "text": (
                    "Describe this figure, chart, graph, diagram, or image in detail "
                    "for a research document context. "
                    "If it contains data or statistics, extract key values, trends, and labels. "
                    "If it contains text, transcribe it. Be specific and thorough."
                )},
            ],
        }],
        max_tokens=500,
    )
    return resp.choices[0].message.content


async def _upload_image(storage, job_id: str, xref: int, ext: str, image_bytes: bytes) -> Optional[str]:
    """Upload image bytes to the Supabase 'images' bucket and return a public URL."""
    path = f"{job_id}/{xref}.{ext}"
    mime = "image/jpeg" if ext in ("jpg", "jpeg") else f"image/{ext}"
    loop = asyncio.get_event_loop()
    try:
        await loop.run_in_executor(
            None,
            lambda: storage.from_("images").upload(
                path, image_bytes, {"content-type": mime, "upsert": "true"}
            ),
        )
        return storage.from_("images").get_public_url(path)
    except Exception:
        return None


async def extract_pdf_multimodal(
    pdf_bytes: bytes,
    storage,
    job_id: str,
) -> list[dict]:
    """
    Extract all content from a PDF as a flat list of chunk dicts:
        {"text": str, "chunk_type": "text"|"image_description", "image_url": str|None}

    Native PDFs: text is chunked normally; embedded images are described by the
    vision model and stored in Supabase.

    Scanned PDFs (little/no extractable text): each page is rendered at 150 DPI
    and described by the vision model (page-level fallback).
    """
    client = get_openai_client()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    result: list[dict] = []

    if _is_scanned_pdf(doc):
        for page in doc:
            try:
                mat = fitz.Matrix(150 / 72, 150 / 72)
                pixmap = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
                jpeg_bytes = pixmap.tobytes("jpeg")
                desc = await asyncio.to_thread(
                    _describe_image_vision_sync, jpeg_bytes, "jpeg", client
                )
                result.append({
                    "text": f"[Page {page.number + 1}] {desc}",
                    "chunk_type": "image_description",
                    "image_url": None,
                })
            except Exception:
                continue
    else:
        non_table_parts: list[str] = []

        for page in doc:
            try:
                tabs = page.find_tables()
                table_rects = []
                for tab in tabs.tables:
                    md = _table_to_markdown(tab.extract())
                    if not md:
                        continue
                    table_rects.append(tab.bbox)
                    table_text = f"[Table, Page {page.number + 1}]\n{md}"
                    table_chunks = chunk_text(table_text) if len(table_text) > CHUNK_SIZE else [table_text]
                    for chunk in table_chunks:
                        result.append({"text": chunk, "chunk_type": "text", "image_url": None})
                non_table_parts.append(_page_text_excluding_tables(page, table_rects))
            except Exception:
                non_table_parts.append(page.get_text())

        for chunk in chunk_text("\n\n".join(non_table_parts)):
            result.append({"text": chunk, "chunk_type": "text", "image_url": None})

        seen_xrefs: set[int] = set()
        for page in doc:
            for img_meta in page.get_images(full=True):
                xref = img_meta[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)

                try:
                    base_image = doc.extract_image(xref)
                except Exception:
                    continue

                img_bytes = base_image["image"]
                ext = base_image["ext"]

                if len(img_bytes) < MIN_IMAGE_BYTES:
                    continue

                image_url = await _upload_image(storage, job_id, xref, ext, img_bytes)

                try:
                    desc = await asyncio.to_thread(
                        _describe_image_vision_sync, img_bytes, ext, client
                    )
                except Exception:
                    continue

                result.append({
                    "text": desc,
                    "chunk_type": "image_description",
                    "image_url": image_url,
                })

    return result


def _tabular_row_chunks(
    filename: str,
    sheet_label: Optional[str],
    columns: list[str],
    rows: list[dict],
) -> list[dict]:
    """
    Serialize tabular rows into RAG-ready chunks where each chunk carries full
    column context so it is independently meaningful when retrieved.

    Each row is serialized as "Row N: Col: val, Col: val, ..." and batches are
    kept under CHUNK_SIZE so the LLM receives a coherent slice of the table.
    """
    if not rows:
        return []

    sheet_info = f" | Sheet: {sheet_label}" if sheet_label else ""
    header = f"[{filename}{sheet_info} | Columns: {', '.join(columns)}]"

    result: list[dict] = []
    current_rows: list[str] = []
    current_len = len(header) + 1

    for i, row in enumerate(rows, 1):
        parts = [f"{col}: {row.get(col)}" for col in columns
                 if row.get(col) is not None and str(row.get(col)).strip()]
        if not parts:
            continue
        row_text = f"Row {i}: " + ", ".join(parts)

        if current_rows and current_len + len(row_text) + 1 > CHUNK_SIZE:
            result.append({
                "text": header + "\n" + "\n".join(current_rows),
                "chunk_type": "text",
                "image_url": None,
            })
            current_rows = []
            current_len = len(header) + 1

        current_rows.append(row_text)
        current_len += len(row_text) + 1

    if current_rows:
        result.append({
            "text": header + "\n" + "\n".join(current_rows),
            "chunk_type": "text",
            "image_url": None,
        })

    return result


def extract_csv(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract CSV data into RAG chunks with a schema summary + row-level chunks."""
    import csv
    import io

    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1")

    reader = csv.DictReader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return []

    columns = [c for c in (reader.fieldnames or []) if c is not None]
    schema_text = (
        f"File: {filename}\n"
        f"Type: CSV spreadsheet\n"
        f"Rows: {len(rows)}\n"
        f"Columns ({len(columns)}): {', '.join(columns)}"
    )
    result: list[dict] = [{"text": schema_text, "chunk_type": "text", "image_url": None}]
    result.extend(_tabular_row_chunks(filename, None, columns, rows))
    return result


def extract_excel(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract all sheets from an XLSX file into RAG chunks."""
    import io
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    result: list[dict] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows = [
            r for r in ws.iter_rows(values_only=True)
            if any(c is not None and str(c).strip() for c in r)
        ]

        if len(all_rows) < 2:
            continue

        columns = [
            str(h).strip() if h is not None else f"Column{i + 1}"
            for i, h in enumerate(all_rows[0])
        ]
        data_rows = [
            {col: val for col, val in zip(columns, row)}
            for row in all_rows[1:]
        ]

        schema_text = (
            f"File: {filename} | Sheet: {sheet_name}\n"
            f"Type: Excel spreadsheet\n"
            f"Rows: {len(data_rows)}\n"
            f"Columns ({len(columns)}): {', '.join(columns)}"
        )
        result.append({"text": schema_text, "chunk_type": "text", "image_url": None})
        result.extend(_tabular_row_chunks(filename, sheet_name, columns, data_rows))

    wb.close()
    return result


def extract_docx(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract Word document content preserving headings and tables."""
    import io
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(file_bytes))
    parts: list[str] = []

    for child in doc.element.body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag == "p":
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            parts.append(f"## {text}" if "Heading" in style else text)
        elif tag == "tbl":
            table = Table(child, doc)
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            md = _table_to_markdown(rows)
            if md:
                parts.append(md)

    full_text = "\n\n".join(parts)
    if not full_text.strip():
        return []

    result: list[dict] = [{"text": f"File: {filename}\nType: Word document", "chunk_type": "text", "image_url": None}]
    for chunk in chunk_text(full_text):
        result.append({"text": chunk, "chunk_type": "text", "image_url": None})
    return result


def extract_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract PowerPoint slides as individual chunks, including speaker notes and tables."""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    result: list[dict] = [{
        "text": f"File: {filename}\nType: PowerPoint presentation\nSlides: {len(prs.slides)}",
        "chunk_type": "text",
        "image_url": None,
    }]

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"[Slide {i}]"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                md = _table_to_markdown(rows)
                if md:
                    parts.append(md)

        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker notes: {notes}]")
        except Exception:
            pass

        slide_text = "\n".join(parts)
        if len(slide_text) <= len(f"[Slide {i}]"):
            continue

        if len(slide_text) > CHUNK_SIZE:
            for chunk in chunk_text(slide_text):
                result.append({"text": chunk, "chunk_type": "text", "image_url": None})
        else:
            result.append({"text": slide_text, "chunk_type": "text", "image_url": None})

    return result


def extract_text_file(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract plain text and markdown files."""
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1")

    if not text.strip():
        return []

    result: list[dict] = []
    for chunk in chunk_text(text):
        result.append({"text": chunk, "chunk_type": "text", "image_url": None})
    return result


def chunk_text(text):

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP
    )

    return splitter.split_text(text)


def raise_openai_http_error(error: OpenAIError):
    error_body = getattr(error, "body", None) or {}
    if not isinstance(error_body, dict):
        error_body = {}

    error_code = getattr(error, "code", None) or error_body.get("code")

    if isinstance(error, AuthenticationError):
        status_code = 401
        detail = "OpenAI authentication failed. Check OPENAI_API_KEY."
    elif isinstance(error, RateLimitError) and error_code == "insufficient_quota":
        status_code = 429
        detail = (
            "OpenAI quota exceeded. Check your OpenAI plan, billing, or usage limits."
        )
    elif isinstance(error, RateLimitError):
        status_code = 429
        detail = "OpenAI rate limit reached. Try again later."
    else:
        status_code = 502
        detail = "OpenAI request failed. Try again later."

    raise HTTPException(status_code=status_code, detail=detail) from error


def raise_pinecone_http_error(error: PineconeApiException):
    error_text = str(error)

    if "Vector dimension" in error_text and "does not match" in error_text:
        detail = (
            "Pinecone vector dimension mismatch. "
            f"Your app is configured for {EMBEDDING_DIMENSIONS}-dimensional "
            "embeddings; make sure the Pinecone index uses the same dimension."
        )
        raise HTTPException(status_code=400, detail=detail) from error

    raise HTTPException(
        status_code=502,
        detail="Pinecone request failed. Try again later."
    ) from error


def store_chunks(chunks, fileName, namespace: str = ""):
    index = get_index()

    embeddings = OpenAIEmbeddings(
        model=EMBEDDING_MODEL,
        dimensions=EMBEDDING_DIMENSIONS
    )

    try:
        chunk_embeddings = embeddings.embed_documents(chunks)
    except OpenAIError as error:
        raise_openai_http_error(error)

    safe_name = fileName.replace("/", "_").replace(" ", "_")
    prefix = f"{namespace}-{safe_name}" if namespace else safe_name
    vectors = []
    for i, (chunk, embedding) in enumerate(zip(chunks, chunk_embeddings)):
        vectors.append({
            "id": f"{prefix}-chunk-{i}",
            "values": embedding,
            "metadata": {
                "text": chunk,
                "source": fileName,
                "chunk_index": i
            }
        })

    kwargs = {"vectors": vectors}
    if namespace:
        kwargs["namespace"] = namespace

    try:
        index.upsert(**kwargs)
    except PineconeApiException as error:
        raise_pinecone_http_error(error)

def ask_question(question: str, namespace: str = ""):
    index = get_index()

    embeddings = OpenAIEmbeddings(
        model=EMBEDDING_MODEL,
        dimensions=EMBEDDING_DIMENSIONS
    )

    try:
        question_embedding = embeddings.embed_query(question)
    except OpenAIError as error:
        raise_openai_http_error(error)

    query_kwargs = {
        "vector": question_embedding,
        "top_k": TOP_K,
        "include_metadata": True,
    }
    if namespace:
        query_kwargs["namespace"] = namespace

    try:
        results = index.query(**query_kwargs)
    except PineconeApiException as error:
        raise_pinecone_http_error(error)

    context = "\n\n".join(
        match["metadata"]["text"]
        for match in results["matches"]
    )

    client = get_openai_client()
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Answer the user's question only "
                        "using the provided context."
                    )
                },
                {
                    "role": "user",
                    "content": f"""
                    Context:
                    {context}

                    Question:
                    {question}
                    """
                }
            ]
        )
    except OpenAIError as error:
        raise_openai_http_error(error)
    
    return {
        "answer": response.choices[0].message.content,
        "sources": [
            {
                "source": match["metadata"].get("source"),
                "chunk_index": match["metadata"].get("chunk_index"),
                "text": match["metadata"]["text"],
                "score": match["score"]
            }
        for match in results["matches"]
        ]
    }
